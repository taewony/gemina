from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os
import re

# --- 설정 ---
MARKDOWN_FILE = "GEMINI.md"
OUTPUT_FILE = "presentation_from_mardown.pptx"
# --------------

def parse_markdown_slides(markdown_content):
    """마크다운을 슬라이드 단위로 분리합니다."""
    # 마크다운 파일 마지막의 연속된 개행 문자를 제거하여 빈 슬라이드가 생기는 것을 방지합니다.
    return re.split(r'\n---\n', markdown_content.strip())

def parse_slide_content(slide_md):
    """한 슬라이드의 마크다운에서 제목, 일반 텍스트, 코드 블록을 분리합니다."""
    title = ""
    normal_content = []
    code_blocks = []
    
    lines = slide_md.strip().split('\n')
    in_code_block = False
    current_code_block = []

    # 제목 추출 로직 개선
    if lines and (lines[0].strip().startswith('# ') or lines[0].strip().startswith('## ')):
        title = lines[0].lstrip('# ').strip()
        lines = lines[1:]
    elif lines:
        # #, ##가 없으면 첫 줄을 제목으로 사용하지 않고, 일반 콘텐츠로 처리할 수 있습니다.
        pass

    for line in lines:
        if line.strip() == '```':
            if in_code_block:
                code_blocks.append('\n'.join(current_code_block))
                current_code_block = []
            in_code_block = not in_code_block
        elif in_code_block:
            current_code_block.append(line)
        else:
            normal_content.append(line)
            
    return title, '\n'.join(normal_content).strip(), code_blocks

def add_normal_content(text_frame, content):
    """일반 텍스트를 파싱하여 글머리 기호와 굵은 글씨를 적용합니다."""
    text_frame.clear()
    lines = content.split('\n')
    
    if not lines or (len(lines) == 1 and not lines[0]):
        return

    p = text_frame.paragraphs[0]
    for i, line in enumerate(lines):
        if i > 0:
            p = text_frame.add_paragraph()
        
        stripped_line = line.lstrip(' ')
        p.level = (len(line) - len(stripped_line)) // 2
        
        # [수정] 글머리 기호 문자를 제거하는 로직을 삭제하고, 원본 텍스트(stripped_line)를 그대로 사용합니다.
        # 이렇게 하면 마크다운의 '-', '*' 등의 문자가 슬라이드에 그대로 표시됩니다.
        line_content = stripped_line
        
        parts = re.split(r'(\*\*.+?\*\*)', line_content)
        for part in filter(None, parts):
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run.text = part
            run.font.size = Pt(16)

def add_code_block(slide, code_content, top_position):
    """슬라이드의 지정된 위치에 코드 블록을 추가합니다."""
    # [원복] 슬라이드 크기에 맞춰 너비 조정 (4:3 비율 기준)
    left = Inches(0.5)
    width = Inches(9) 
    
    num_lines = len(code_content.split('\n')) + 1
    height = Pt(12 * 1.5 * num_lines)
    
    max_height = Inches(7.5) - top_position - Inches(0.5)
    if height > max_height:
        height = max_height

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top_position, width, height)
    
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    line = shape.line
    line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.text = code_content
    p.alignment = PP_ALIGN.LEFT
    
    # [핵심 수정 유지] 단락의 위/아래 공백을 0으로 설정
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    font = p.font
    font.name = 'Courier New'
    font.size = Pt(12)
    font.color.rgb = RGBColor(50, 50, 50)
    
    return shape

def create_presentation(slides_raw):
    """파싱된 데이터로 PPTX 프레젠테이션을 생성합니다."""
    prs = Presentation()
    # [원복] 슬라이드 크기 설정 제거 (기본 4:3 사용)
    
    for slide_md in slides_raw:
        title, normal_content, code_blocks = parse_slide_content(slide_md)
        
        slide_layout = prs.slide_layouts[5] 
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title

        current_top = title_shape.top + title_shape.height

        if normal_content:
            # [원복] 너비를 4:3 비율에 맞게 조정
            left = Inches(0.5)
            width = Inches(9)
            height = Inches(1) 
            
            body_shape = slide.shapes.add_textbox(left, current_top, width, height)
            add_normal_content(body_shape.text_frame, normal_content)
            
            current_top = body_shape.top + body_shape.height + Inches(0.2)

        if code_blocks:
            for code in code_blocks:
                code_shape = add_code_block(slide, code, current_top)
                current_top = code_shape.top + code_shape.height + Inches(0.1)

    return prs

def main():
    print(f"'{MARKDOWN_FILE}' 파일을 읽어 PPTX 파일을 생성합니다.")
    # [원복] 파일 경로 처리 방식을 원래 코드로 복원
    try:
        with open(os.path.join("..", MARKDOWN_FILE), 'r', encoding='utf-8') as f:
            markdown_content = f.read()
    except FileNotFoundError:
        print(f"오류: '{MARKDOWN_FILE}'을 찾을 수 없습니다. (스크립트 상위 폴더에 있어야 합니다)")
        return

    slides_raw = parse_markdown_slides(markdown_content)
    prs = create_presentation(slides_raw)
    
    # [원복] 파일 저장 경로를 원래 코드로 복원
    prs.save(OUTPUT_FILE)
    print(f"'{OUTPUT_FILE}' 파일이 성공적으로 개선되어 생성되었습니다.")

if __name__ == '__main__':
    # [원복] 스크립트 실행 경로 변경 로직 복원
    abspath = os.path.abspath(__file__)
    dname = os.path.dirname(abspath)
    os.chdir(dname)
    main()
